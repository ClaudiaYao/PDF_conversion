import os
import json
import pandas as pd
from langchain.prompts import PromptTemplate
from langchain.output_parsers import PydanticOutputParser
from pydantic import BaseModel, Field, validator
from langchain_openai import ChatOpenAI
import pptx
from typing import List, Optional, Literal
from dotenv import load_dotenv
from PIL import Image
from transformers import BlipProcessor, BlipForConditionalGeneration


load_dotenv()
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")

# For image captioning
blip_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-large")
blip_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-large")

# Json structure to parse
class Paragraph(BaseModel):
    text: str = Field(description="text content of the paragraph", default='')
    indent_level: Literal[0, 1, 2] = Field(description="indent level of the text", default=0)

# For section level generation
class ContentSlide(BaseModel):
    # slide_layout: Literal[0, 1] = Field(description=slide_layout_descriptions)
    title: str = Field(description="title of the slide")
    content: List[Paragraph] = Field(description="list of content paragraphs", default=[])

class ContentSlides(BaseModel):
    slides: List[ContentSlide] = Field(description="list of content slides")

class TitleSlide(BaseModel):
    # slide_layout: Literal[0, 1] = Field(description=slide_layout_descriptions)
    title: str = Field(description="title of the slide")
    subtitle: str = Field(description="subtitle of the slide", default='')

# For full summary generation
# From https://python-pptx.readthedocs.io/en/latest/user/slides.html
slide_layout_descriptions = {0: "Title (presentation title slide)", 1: "Title and Content"}

class Slide(BaseModel):
    slide_layout: Literal[0, 1] = Field(description=slide_layout_descriptions)
    title: str = Field(description="title of the slide")
    subtitle: str = Field(description="subtitle of the slide, only for title slide", default='')
    content: List[Paragraph] = Field(description="list of content paragraphs", default=[])

class Presentation(BaseModel):
    slides: List[Slide] = Field(description="list of slides")

def get_meta(article_name):
    df_meta = pd.read_csv(f'processed/{article_name}_meta.csv')
    paper_title = df_meta['Title'][0]
    abstract = df_meta['Abstract'][0]

    return paper_title, abstract

def get_section_groundtruth(section):
    section_summary = ""
    name = section['Section']
    text = section['Text']
    subsections = section['Subsections']
    groundtruth = section['Groundtruth']

    if name != "NA":
        section_summary += name + '\n'
    if groundtruth:
        section_summary += groundtruth + '\n'
    
    if subsections:
        for subsection in subsections:
            section_summary += get_section_groundtruth(subsection)

    return section_summary

def get_section_summary(section):
    with open(f'model_summarizer/results/model-summary_results.json', encoding='utf-8') as f:
        data = json.load(f)
    def build_dict(seq, key):
        return dict((d[key], dict(d, index=index)) for (index, d) in enumerate(seq))

    summary_by_section_name = build_dict(data, key="Section Name")
    section_summary = ""
    name = section['Section']
    text = section['Text']
    subsections = section['Subsections']
    if name not in summary_by_section_name:
        return ""
    summary = summary_by_section_name.get(name)["Generated Summary"]

    if name != "NA":
        section_summary += name + '\n'
    if summary:
        section_summary += summary + '\n'
    
    if subsections:
        for subsection in subsections:
            section_summary += get_section_summary(subsection)

    return section_summary

def get_section_image_paths(section):
    processed_folder = "processed"

    image_paths = []
    section_num = section["Section_Num"]


    for filename in os.listdir(processed_folder):
        if filename.startswith(f'section_{section_num}_') or filename.startswith(f'section_{section_num}-'):
            raw_image = Image.open(f'processed/{filename}').convert('RGB')

            # unconditional image captioning
            inputs = blip_processor(raw_image, return_tensors="pt")

            out = blip_model.generate(**inputs)
            caption = blip_processor.decode(out[0], skip_special_tokens=True)
            image_paths.append({'path': filename, 'caption': caption})
    return image_paths

def get_title_slide_data(article_name):
    llm = ChatOpenAI(model="gpt-3.5-turbo-0125", temperature=0, openai_api_key=OPENAI_API_KEY)

    paper_title, abstract = get_meta(article_name)

    title_slide_parser = PydanticOutputParser(pydantic_object=TitleSlide)

    title_slide_prompt = PromptTemplate(
        template="Answer the user query.\n{format_instructions}\n{query}\n",
        input_variables=["query"],
        partial_variables={"format_instructions": title_slide_parser.get_format_instructions()},
    )

    title_slide_query = f"""
Generate a title slide for a technical presentation based on the title and abstract of a paper. Subtitle should not be too long
Paper title:
{paper_title}
Abstract:
{abstract}
"""
    title_slide_chain = title_slide_prompt | llm | title_slide_parser
    title_slide_response = title_slide_chain.invoke({"query": title_slide_query})
    title_slide_data = json.loads(title_slide_response.model_dump_json())
    
    return title_slide_data

def get_toc_slide_data(article_name, section_names):
    llm = ChatOpenAI(model="gpt-3.5-turbo-0125", temperature=0, openai_api_key=OPENAI_API_KEY)

    paper_title, abstract = get_meta(article_name)

    toc_slide_parser = PydanticOutputParser(pydantic_object=ContentSlide)

    toc_slide_prompt = PromptTemplate(
        template="Answer the user query.\n{format_instructions}\n{query}\n",
        input_variables=["query"],
        partial_variables={"format_instructions": toc_slide_parser.get_format_instructions()},
    )

    toc_slide_query = f"""
Generate a table of contents slide for a technical presentation based on the title and section names of a paper. Only include relevant section names (ie. skip no titles)
Paper title:
{paper_title}
Section names:
{", ".join(section_names)}
"""
    toc_slide_chain = toc_slide_prompt | llm | toc_slide_parser
    toc_slide_response = toc_slide_chain.invoke({"query": toc_slide_query})
    toc_slide_data = json.loads(toc_slide_response.model_dump_json())
    
    return toc_slide_data

def get_content_slide_datas(section_names, section_texts):
    llm = ChatOpenAI(model="gpt-3.5-turbo-0125", temperature=0, openai_api_key=OPENAI_API_KEY)

    content_slide_parser = PydanticOutputParser(pydantic_object=ContentSlides)

    content_slide_prompt = PromptTemplate(
        template="Answer the user query.\n{format_instructions}\n{query}\n",
        input_variables=["query"],
        partial_variables={"format_instructions": content_slide_parser.get_format_instructions()},
    )

    content_slide_chain = content_slide_prompt | llm | content_slide_parser
    content_slide_datas = []
    completed_sections = []

    for i, section_text in enumerate(section_texts):
        section_name = section_names[i]

        content_slide_query = f"""
Generate a portion of a technical presentation for a section of a paper based on its summary. There should not be too much text per slide. Always give a relevant title.
Section summary:
{section_text}

"""
        
        if completed_sections:
            content_slide_query += f"""
The previous slides have covered the following sections, make sure the slides are coherent with the previous slides (eg. do not generate any introduction slides):
{", ".join(completed_sections)}

"""

        content_slide_response = content_slide_chain.invoke({"query": content_slide_query})
        r = json.loads(content_slide_response.model_dump_json())
        slide_data = r["slides"]
        content_slide_datas.append(slide_data)
        completed_sections.append(section_name)
    
    return content_slide_datas

def get_full_summary_data(full_summary):
    full_summary_parser = PydanticOutputParser(pydantic_object=Presentation)

    full_summary_prompt = PromptTemplate(
        template="Answer the user query.\n{format_instructions}\n{query}\n",
        input_variables=["query"],
        partial_variables={"format_instructions": full_summary_parser.get_format_instructions()},
    )

    full_summary_query = f"""
Generate a technical presentation based on the section-by-section summary of a paper below. The first slide should be a title only slide. The second slide should be a table of contents slide
Paper summary:
{full_summary}
"""
    full_summary_chain = full_summary_prompt | llm | full_summary_parser
    full_summary_response = full_summary_chain.invoke({"query": full_summary_query})
    r = json.loads(full_summary_response.model_dump_json())
    slide_data = r["slides"]

    return slide_data

def add_image(slide, placeholder_id, image_url):
    placeholder = slide.placeholders[placeholder_id]
    # Calculate the image size of the image
    im = Image.open(image_url)
    width, height = im.size
    # Make sure the placeholder doesn't zoom in
    placeholder.height = height
    placeholder.width = width
    # Insert the picture
    placeholder = placeholder.insert_picture(image_url)
    # Calculate ratios and compare
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio
    # Placeholder width too wide:
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    # Placeholder height too high
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side

def generate_section_level_ppt(theme, title_slide_data, toc_slide_data, content_slide_datas, section_image_paths):
    prs = pptx.Presentation(f"data/ppt_themes/{theme}.pptx")

    # Add title slide
    slide_layout = prs.slide_layouts[0] # title slide layout
    new_slide = prs.slides.add_slide(slide_layout)
    shapes = new_slide.shapes
    body_shape = shapes.placeholders[0]
    tf = body_shape.text_frame
    if title_slide_data['title']:
        title = new_slide.shapes.title
        title.text = title_slide_data['title']
        if theme == "Gallery":
            tf.fit_text(max_size=40, bold=True)

    if title_slide_data['subtitle']:
        subtitle = new_slide.placeholders[1]
        subtitle.text = title_slide_data['subtitle']

    # Add toc slide
    slide_layout = prs.slide_layouts[1] # content slide layout
    new_slide = prs.slides.add_slide(slide_layout)
    shapes = new_slide.shapes
    body_shape = shapes.placeholders[0]
    tf = body_shape.text_frame
    if toc_slide_data['title']:
        title = new_slide.shapes.title
        title.text = toc_slide_data['title']
        
    if toc_slide_data['content']:
        shapes = new_slide.shapes
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        for i, paragraph in enumerate(toc_slide_data['content']):
            if i == 0:
                tf.text = paragraph['text']
                # tf.fit_text(font_family="Calibri", max_size=18, bold=True)
            else:
                p = tf.add_paragraph()
                p.text = paragraph['text']
                p.level = paragraph['indent_level']
                # tf.fit_text(font_family="Calibri", max_size=18, bold=True)

    for i, slide in enumerate(content_slide_datas):
        image_paths = section_image_paths[i]
        slide = slide[0]
        slide_layout = prs.slide_layouts[1] # Content slide
        new_slide = prs.slides.add_slide(slide_layout)
        # for shape in new_slide.placeholders:
        #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
        
        if slide['title']:
            title = new_slide.shapes.title
            title.text = slide['title']
        
        if slide['content']:
            shapes = new_slide.shapes
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            for i, paragraph in enumerate(slide['content']):
                if i == 0:
                    tf.text = paragraph['text']
                    # tf.fit_text(font_family="Calibri", max_size=18, bold=True)
                else:
                    p = tf.add_paragraph()
                    p.text = paragraph['text']
                    p.level = paragraph['indent_level']
                    # tf.fit_text(font_family="Calibri", max_size=18, bold=True)
        if image_paths:
                for image_path in image_paths:
                    slide_layout = prs.slide_layouts[8] # image_with_caption slide
                    new_slide = prs.slides.add_slide(slide_layout)
                    shapes = new_slide.shapes
                    body_shape = shapes.placeholders[0]
                    tf = body_shape.text_frame
                    tf.text = image_path['caption']
                    add_image(new_slide,1,f"processed/{image_path['path']}")

    return prs

def generate_full_summary_ppt(theme, full_summary_data):
    prs = pptx.Presentation(f"ppt_themes/{theme}.pptx")

    for slide in full_summary_data:
        slide_layout = prs.slide_layouts[slide['slide_layout']]
        new_slide = prs.slides.add_slide(slide_layout)
        for shape in new_slide.placeholders:
            print('%d %s' % (shape.placeholder_format.idx, shape.name))
        
        if slide['title']:
            title = new_slide.shapes.title
            title.text = slide['title']

        if slide['slide_layout'] == 0 and slide['subtitle']:
            subtitle = new_slide.placeholders[1]
            subtitle.text = slide['subtitle']
        
        if slide['slide_layout'] == 1 and slide['content']:
            shapes = new_slide.shapes
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            for i, paragraph in enumerate(slide['content']):
                if i == 0:
                    tf.text = paragraph['text']
                    tf.fit_text(font_family="Calibri", max_size=18, bold=True)
                else:
                    p = tf.add_paragraph()
                    p.text = paragraph['text']
                    p.level = paragraph['indent_level']
                    tf.fit_text(font_family="Calibri", max_size=18, bold=True)

    return prs
                    
